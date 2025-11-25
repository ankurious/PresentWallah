from docx import Document
from docx.shared import Pt, Inches, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import List
from app.models import Project, Section
from app.services.image import image_service
import io

class DocumentService:
    
    # Template color schemes
    TEMPLATES = {
        "modern": {
            "primary": RGBColor(15, 32, 96),      # Navy blue
            "secondary": RGBColor(30, 58, 138),   # Lighter blue
            "accent": RGBColor(220, 230, 255),    # Light blue
            "text": RGBColor(40, 40, 40)
        },
        "minimal": {
            "primary": RGBColor(50, 50, 50),      # Dark gray
            "secondary": RGBColor(100, 100, 100), # Medium gray
            "accent": RGBColor(240, 240, 240),    # Light gray
            "text": RGBColor(60, 60, 60)
        },
        "corporate": {
            "primary": RGBColor(0, 51, 102),      # Corporate blue
            "secondary": RGBColor(0, 102, 204),   # Bright blue
            "accent": RGBColor(230, 240, 255),    # Very light blue
            "text": RGBColor(30, 30, 30)
        },
        "creative": {
            "primary": RGBColor(123, 31, 162),    # Purple
            "secondary": RGBColor(156, 39, 176),  # Lighter purple
            "accent": RGBColor(243, 229, 245),    # Light purple
            "text": RGBColor(50, 50, 50)
        }
    }
    
    @staticmethod
    def get_template_colors(template_name: str):
        """Get color scheme for a template"""
        return DocumentService.TEMPLATES.get(template_name, DocumentService.TEMPLATES["modern"])
    @staticmethod
    def generate_docx(project: Project, sections: List[Section]) -> io.BytesIO:
        """Generate a Word document from project sections"""
        doc = Document()
        
        # Add title
        title = doc.add_heading(project.title, 0)
        title.alignment = 1  # Center alignment
        
        # Add topic/description
        doc.add_paragraph(f"Topic: {project.main_topic}")
        doc.add_paragraph()  # Empty line
        
        # Sort sections by order
        sorted_sections = sorted(sections, key=lambda s: s.order)
        
        # Add each section
        for section in sorted_sections:
            # Section heading
            doc.add_heading(section.title, 1)
            
            # Section content
            if section.content:
                doc.add_paragraph(section.content)
            else:
                doc.add_paragraph("[Content not generated yet]")
            
            # Add space between sections
            doc.add_paragraph()
        
        # Save to BytesIO
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream
    
    @staticmethod
    def generate_pptx(project: Project, sections: List[Section], include_images: bool = True) -> io.BytesIO:
        """Generate a PowerPoint presentation with professional templates and images"""
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Get template colors
        colors = DocumentService.get_template_colors(project.template if hasattr(project, 'template') else "modern")
        font_size = project.font_size if hasattr(project, 'font_size') and project.font_size else 20
        
        # Sort sections by order
        sorted_sections = sorted(sections, key=lambda s: s.order)
        
        for idx, section in enumerate(sorted_sections):
            if idx == 0:
                # First slide - Professional Title Slide
                DocumentService._create_title_slide(prs, section, project, include_images, colors)
            else:
                # Content slides with images
                DocumentService._create_content_slide(prs, section, include_images, colors, font_size)
        
        # Save to BytesIO
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream
    
    @staticmethod
    def _create_title_slide(prs: Presentation, section: Section, project: Project, include_images: bool, colors: dict):
        """Create a modern, professional title slide with gradient effect"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Gradient background simulation with shapes
        bg1 = slide.shapes.add_shape(
            1, PptxInches(0), PptxInches(0),
            PptxInches(10), PptxInches(7.5)
        )
        bg1.fill.solid()
        bg1.fill.fore_color.rgb = colors["primary"]
        bg1.line.fill.background()
        
        # Accent shape for visual interest
        accent = slide.shapes.add_shape(
            1, PptxInches(0), PptxInches(5),
            PptxInches(10), PptxInches(2.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = colors["secondary"]
        accent.line.fill.background()
        
        # Main title
        title_box = slide.shapes.add_textbox(
            PptxInches(0.8), PptxInches(2.5),
            PptxInches(8.4), PptxInches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = section.title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.name = "Calibri"
        title_para.font.size = PptxPt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(5.5),
            PptxInches(8), PptxInches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = project.main_topic
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.name = "Calibri"
        subtitle_para.font.size = PptxPt(28)
        subtitle_para.font.color.rgb = colors["accent"]
    
    @staticmethod
    def _create_content_slide(prs: Presentation, section: Section, include_images: bool, colors: dict, font_size: int = 20):
        """Create modern content slide with image and styled bullets"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # White background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Modern header bar with template color
        header = slide.shapes.add_shape(
            1, PptxInches(0), PptxInches(0),
            PptxInches(10), PptxInches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors["primary"]
        header.line.fill.background()
        
        # Title in header
        title_box = slide.shapes.add_textbox(
            PptxInches(0.6), PptxInches(0.25),
            PptxInches(8.8), PptxInches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.text = section.title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Calibri"
        title_para.font.size = PptxPt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Try to fetch and add image
        image_added = False
        if include_images:
            # Try multiple search terms for better image results
            search_queries = [
                section.title,
                f"{section.title} business",
                f"{section.title} professional"
            ]
            
            for query in search_queries:
                image_url = image_service.search_image(query)
                if image_url:
                    try:
                        image_data = image_service.download_image(image_url)
                        if image_data:
                            image_stream = io.BytesIO(image_data)
                            # Add image on the left with rounded effect
                            pic = slide.shapes.add_picture(
                                image_stream,
                                PptxInches(0.5), PptxInches(1.7),
                                width=PptxInches(4.2), height=PptxInches(4.8)
                            )
                            image_added = True
                            print(f"Image added for: {section.title}")
                            break
                    except Exception as e:
                        print(f"Error adding image for '{query}': {e}")
                        continue
        
        # Position content based on whether image was added
        if image_added:
            content_left = PptxInches(5.2)
            content_width = PptxInches(4.3)
        else:
            content_left = PptxInches(1)
            content_width = PptxInches(8)
        
        # Content area with professional bullet formatting
        content_box = slide.shapes.add_textbox(
            content_left, PptxInches(1.8),
            content_width, PptxInches(4.8)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        if section.content:
            # Parse content into bullet points
            lines = [line.strip() for line in section.content.strip().split('\n') if line.strip()]
            
            for i, line in enumerate(lines):
                # Remove any existing bullet characters
                line = line.lstrip('•-*►▪').strip()
                if not line:
                    continue
                
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line
                p.level = 0
                p.font.name = "Calibri"
                p.font.size = PptxPt(font_size)
                p.font.color.rgb = RGBColor(40, 40, 40)
                p.space_before = PptxPt(14)
                p.space_after = PptxPt(8)
                p.line_spacing = 1.3
                
                # Add bullet point character
                if i > 0 or len(lines) > 1:
                    p.text = "• " + line
        else:
            p = text_frame.paragraphs[0]
            p.text = "Content not generated yet"
            p.font.name = "Calibri"
            p.font.size = PptxPt(20)
            p.font.italic = True
            p.font.color.rgb = RGBColor(150, 150, 150)

document_service = DocumentService()
